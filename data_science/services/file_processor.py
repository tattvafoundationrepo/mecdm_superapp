"""File upload and text extraction service for GCS-backed document processing."""

import io
import logging
import os
import re
import uuid

from google.cloud import storage

logger = logging.getLogger(__name__)

UPLOAD_BUCKET_NAME = os.environ.get("UPLOAD_BUCKET_NAME", "mecdm-superapp-uploads")

# Mime type -> { extensions, max_size_bytes, processing_strategy }
ALLOWED_MIME_TYPES: dict[str, dict] = {
    "application/pdf": {
        "extensions": [".pdf"],
        "max_size": 20 * 1024 * 1024,
        "strategy": "gemini_multimodal",
    },
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": {
        "extensions": [".docx"],
        "max_size": 20 * 1024 * 1024,
        "strategy": "text_extraction",
    },
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": {
        "extensions": [".pptx"],
        "max_size": 20 * 1024 * 1024,
        "strategy": "text_extraction",
    },
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": {
        "extensions": [".xlsx"],
        "max_size": 20 * 1024 * 1024,
        "strategy": "text_extraction",
    },
    "image/png": {
        "extensions": [".png"],
        "max_size": 10 * 1024 * 1024,
        "strategy": "gemini_multimodal",
    },
    "image/jpeg": {
        "extensions": [".jpg", ".jpeg"],
        "max_size": 10 * 1024 * 1024,
        "strategy": "gemini_multimodal",
    },
    "image/webp": {
        "extensions": [".webp"],
        "max_size": 10 * 1024 * 1024,
        "strategy": "gemini_multimodal",
    },
}

_gcs_client: storage.Client | None = None


def _get_gcs_client() -> storage.Client:
    """Lazy singleton for the GCS client."""
    global _gcs_client
    if _gcs_client is None:
        _gcs_client = storage.Client()
    return _gcs_client


def sanitize_filename(name: str) -> str:
    """Strip path components, special chars, and limit length."""
    # Take only the basename
    name = os.path.basename(name)
    # Replace anything that isn't alphanumeric, dot, hyphen, or underscore
    name = re.sub(r"[^\w.\-]", "_", name)
    # Collapse multiple underscores
    name = re.sub(r"_+", "_", name)
    # Limit length (preserve extension)
    base, ext = os.path.splitext(name)
    if len(base) > 200:
        base = base[:200]
    return f"{base}{ext}"


def validate_file(filename: str, content_type: str, size: int) -> None:
    """Validate file type, extension, and size. Raises ValueError on failure."""
    if content_type not in ALLOWED_MIME_TYPES:
        allowed = ", ".join(
            ext
            for cfg in ALLOWED_MIME_TYPES.values()
            for ext in cfg["extensions"]
        )
        raise ValueError(
            f"Unsupported file type: {content_type}. Allowed: {allowed}"
        )

    config = ALLOWED_MIME_TYPES[content_type]

    # Check extension matches mime type
    _, ext = os.path.splitext(filename.lower())
    if ext not in config["extensions"]:
        raise ValueError(
            f"File extension '{ext}' does not match content type '{content_type}'. "
            f"Expected: {config['extensions']}"
        )

    if size > config["max_size"]:
        max_mb = config["max_size"] / (1024 * 1024)
        raise ValueError(
            f"File too large ({size / (1024 * 1024):.1f} MB). "
            f"Maximum for {content_type}: {max_mb:.0f} MB"
        )

    if size == 0:
        raise ValueError("File is empty")


def upload_to_gcs(
    file_bytes: bytes,
    filename: str,
    mime_type: str,
    user_id: str,
    session_id: str,
) -> dict:
    """Upload file to GCS and return metadata dict.

    Returns:
        dict with file_id, file_name, mime_type, gcs_uri, size_bytes, processing_strategy
    """
    safe_name = sanitize_filename(filename)
    file_id = str(uuid.uuid4())
    blob_path = f"uploads/{user_id}/{session_id}/{file_id}_{safe_name}"
    gcs_uri = f"gs://{UPLOAD_BUCKET_NAME}/{blob_path}"
    strategy = ALLOWED_MIME_TYPES[mime_type]["strategy"]

    client = _get_gcs_client()
    bucket = client.bucket(UPLOAD_BUCKET_NAME)
    blob = bucket.blob(blob_path)
    blob.upload_from_string(file_bytes, content_type=mime_type)
    logger.info("Uploaded %s (%d bytes) to %s", safe_name, len(file_bytes), gcs_uri)

    # For office documents, extract text and store companion file
    if strategy == "text_extraction":
        try:
            extracted = extract_text(mime_type, file_bytes)
            text_blob = bucket.blob(f"{blob_path}.extracted.txt")
            text_blob.upload_from_string(extracted, content_type="text/plain; charset=utf-8")
            logger.info("Stored extracted text for %s (%d chars)", safe_name, len(extracted))
        except Exception:
            logger.exception("Text extraction failed for %s", safe_name)

    return {
        "file_id": file_id,
        "file_name": safe_name,
        "mime_type": mime_type,
        "gcs_uri": gcs_uri,
        "size_bytes": len(file_bytes),
        "processing_strategy": strategy,
    }


def extract_text(mime_type: str, file_bytes: bytes) -> str:
    """Extract text content from office documents."""
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx(file_bytes)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _extract_pptx(file_bytes)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return _extract_xlsx(file_bytes)
    else:
        raise ValueError(f"No text extractor for mime type: {mime_type}")


def _extract_docx(file_bytes: bytes) -> str:
    """Extract all paragraph text from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    return "\n\n".join(paragraphs)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract text from all slides in a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    sections: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_text.append(" | ".join(cells))
        if len(slide_text) > 1:  # More than just the header
            sections.append("\n".join(slide_text))

    return "\n\n".join(sections)


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract data from XLSX as markdown tables per sheet."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sections: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # Find the first non-empty row as header
        header_idx = 0
        for idx, row in enumerate(rows):
            if any(cell is not None for cell in row):
                header_idx = idx
                break

        header = [str(cell) if cell is not None else "" for cell in rows[header_idx]]
        section_lines = [f"### Sheet: {sheet_name}", "| " + " | ".join(header) + " |"]
        section_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        # Limit to 500 data rows to avoid massive text
        for row in rows[header_idx + 1 : header_idx + 501]:
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c for c in cells):
                section_lines.append("| " + " | ".join(cells) + " |")

        remaining = len(rows) - header_idx - 501
        if remaining > 0:
            section_lines.append(f"\n*... {remaining} more rows truncated*")

        sections.append("\n".join(section_lines))

    wb.close()
    return "\n\n".join(sections)


def read_extracted_text(gcs_uri: str) -> str:
    """Read the extracted text companion file for an uploaded document.

    Args:
        gcs_uri: The GCS URI of the original uploaded file.

    Returns:
        The extracted text content.

    Raises:
        ValueError: If the URI is invalid or doesn't belong to the upload bucket.
        FileNotFoundError: If the extracted text file doesn't exist.
    """
    expected_prefix = f"gs://{UPLOAD_BUCKET_NAME}/"
    if not gcs_uri.startswith(expected_prefix):
        raise ValueError(
            f"Invalid GCS URI. Must start with {expected_prefix}"
        )

    blob_path = gcs_uri[len(expected_prefix):]
    text_blob_path = f"{blob_path}.extracted.txt"

    client = _get_gcs_client()
    bucket = client.bucket(UPLOAD_BUCKET_NAME)
    blob = bucket.blob(text_blob_path)

    if not blob.exists():
        raise FileNotFoundError(
            f"No extracted text found for {gcs_uri}. "
            "The file may not have been processed yet."
        )

    return blob.download_as_text(encoding="utf-8")
